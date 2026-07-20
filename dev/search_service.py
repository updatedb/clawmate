"""
搜索服务层 — 文件名搜索、内容搜索、文本提取。

所有搜索逻辑独立于此模块，不依赖 routes 层。
"""
from __future__ import annotations

import logging
from collections import deque
from pathlib import Path
from typing import Dict, List

from service import safe_path, file_info

logger = logging.getLogger("clawmate.search_service")


def search_media(query: str, root_id: str, rel_dir: str = "", recursive: bool = True,
                 limit: int = 200, max_depth: int = 8, timeout: float = 10.0,
                 exclude_dir: list = None) -> Dict:
    """搜索文件/目录名，支持递归深度限制和硬超时。

    Args:
        max_depth: 递归最大深度（从 start_dir 算起），默认 8 层
        timeout: 搜索超时秒数，默认 10s；超时后返回已有结果并标记 truncated=True
        exclude_dir: 要排除的目录名列表（不区分大小写），排除隐藏目录和临时目录
    """
    import time as _time
    _deadline = _time.time() + timeout

    root_path, target, _ = safe_path(root_id, rel_dir)
    if not target.exists() or not target.is_dir():
        raise FileNotFoundError("Directory not found")

    query_lower = query.lower().strip()
    results: List[Dict] = []
    truncated = False
    if not query_lower:
        return {"query": query, "results": results, "truncated": False}

    exclude_dir_set = set(d.lower() for d in (exclude_dir or []))

    if recursive:
        queue: deque = deque([(target, 0)])  # (dir_path, depth)
        while queue:
            if _time.time() > _deadline:
                truncated = True
                break
            dir_path, depth = queue.popleft()
            if depth > max_depth:
                continue
            try:
                entries = sorted(dir_path.iterdir(), key=lambda p: (not p.is_dir(), p.name.lower()))
            except (OSError, PermissionError):
                continue
            for entry in entries:
                if query_lower in entry.name.lower():
                    results.append(file_info(entry, str(entry.relative_to(root_path))))
                    if len(results) >= limit:
                        break
            if len(results) >= limit:
                break
            if depth < max_depth:
                for entry in entries:
                    if not entry.is_dir():
                        continue
                    # Skip hidden directories
                    if entry.name.startswith("."):
                        continue
                    # Skip excluded directories (e.g. __pycache__, node_modules)
                    if entry.name.lower() in exclude_dir_set:
                        continue
                    queue.append((entry, depth + 1))
    else:
        for entry in target.iterdir():
            if _time.time() > _deadline:
                truncated = True
                break
            if query_lower in entry.name.lower():
                results.append(file_info(entry, str(entry.relative_to(root_path))))
                if len(results) >= limit:
                    break

    return {"query": query, "results": results, "truncated": truncated}


# ── Content search ──────────────────────────────────────────────────────────────

# File extensions that can be content-searched as plain text directly
_CONTENT_TEXT_EXTS = {
    ".txt", ".md", ".markdown", ".mdx", ".json", ".csv", ".log",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".html", ".css", ".yaml", ".yml",
    ".ini", ".toml", ".xml", ".gpx", ".kml", ".conf", ".env", ".sh", ".bat",
    ".ps1", ".sql", ".r", ".go", ".java", ".c", ".cpp", ".h", ".hpp", ".vue",
    ".srt", ".rst", ".tex", ".Makefile", ".dockerfile", ".cfg",
}

# File extensions that need text extraction before search
_CONTENT_EXTRACT_EXTS = {
    ".pdf": "pymupdf",
    ".docx": "python-docx",
    ".xlsx": "openpyxl",
    ".pptx": "python-pptx",
}


_RG_PATH = None

def _check_extractors() -> Dict[str, bool]:
    """Check which document extractors are available.

    Returns a dict mapping extension -> available (bool).
    Python's own import cache makes repeated checks essentially free,
    so no manual cache is needed.
    """
    checks = {
        ".pdf": "fitz",
        ".docx": "docx",
        ".xlsx": "openpyxl",
        ".pptx": "pptx",
    }
    result: Dict[str, bool] = {}
    for ext, mod in checks.items():
        try:
            __import__(mod)
            result[ext] = True
        except ImportError:
            result[ext] = False

    return result


def _rg_available() -> bool:
    """Check if ripgrep is installed."""
    import shutil
    global _RG_PATH
    if _RG_PATH:
        return True
    _RG_PATH = shutil.which("rg")
    if _RG_PATH:
        return True
    # Check common bundled locations (Codex, etc.)
    for candidate in [
        Path.home() / ".npm-global/lib/node_modules/@openai/codex/node_modules/@openai/codex-linux-x64/vendor/x86_64-unknown-linux-musl/codex-path/rg",
        Path.home() / ".npm-global/lib/node_modules/@openai/codex/node_modules/@openai/codex-darwin-x64/vendor/aarch64-apple-darwin/codex-path/rg",
    ]:
        if candidate.is_file():
            _RG_PATH = str(candidate)
            return True
    return False


def _path_hash(file_path: Path) -> str:
    """Short hash of a file path for cache key."""
    import hashlib
    return hashlib.md5(str(file_path).encode()).hexdigest()[:12]


def _is_junk_path(path: Path) -> bool:
    """Check if a path is macOS/editor metadata junk that should be skipped.

    Returns True for:
    - __MACOSX directories and their contents (macOS ZIP resource forks)
    - ._* AppleDouble files (macOS resource fork metadata)
    - Hidden files/directories starting with '.' (except .clawmate/)
    - .DS_Store files
    """
    # Check if any path component is a junk directory or hidden directory
    for part in path.parts:
        if part == "__MACOSX":
            return True
        if part == ".DS_Store":
            return True
        if part.startswith("._"):
            return True
        # Skip hidden directories (except .clawmate which is our own cache)
        if part.startswith(".") and part != ".clawmate":
            return True
    # Also check the filename directly for AppleDouble prefix
    if path.name.startswith("._"):
        return True
    return False


def _extract_by_type(file_path: Path) -> str | None:
    """Extract plain text from a file based on its extension.

    Returns the extracted text, or None if extraction is not supported.
    """
    ext = file_path.suffix.lower()

    if ext == ".pdf":
        try:
            import fitz  # pymupdf
            text_parts = []
            with fitz.open(file_path) as doc:
                for page_num, page in enumerate(doc):
                    page_text = page.get_text()
                    if page_text.strip():
                        text_parts.append(f"--- Page {page_num + 1} ---\n{page_text}")
            return "\n\n".join(text_parts) if text_parts else None
        except ImportError:
            logger.warning("pymupdf not installed, cannot search PDF files")
            return None
        except Exception as e:
            logger.warning(f"Failed to extract text from PDF {file_path}: {e}")
            return None

    if ext == ".docx":
        try:
            from docx import Document
            doc = Document(str(file_path))
            text_parts = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(text_parts) if text_parts else None
        except ImportError:
            logger.warning("python-docx not installed, cannot search DOCX files")
            return None
        except Exception as e:
            logger.warning(f"Failed to extract text from DOCX {file_path}: {e}")
            return None

    if ext == ".xlsx":
        try:
            from openpyxl import load_workbook
            wb = load_workbook(file_path, read_only=True, data_only=True)
            text_parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                sheet_lines = [f"--- Sheet: {sheet_name} ---"]
                for row in ws.iter_rows(values_only=True):
                    row_text = "\t".join(str(c) if c is not None else "" for c in row)
                    if row_text.strip():
                        sheet_lines.append(row_text)
                if len(sheet_lines) > 1:
                    text_parts.append("\n".join(sheet_lines))
            wb.close()
            return "\n\n".join(text_parts) if text_parts else None
        except ImportError:
            logger.warning("openpyxl not installed, cannot search XLSX files")
            return None
        except Exception as e:
            logger.warning(f"Failed to extract text from XLSX {file_path}: {e}")
            return None

    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            text_parts = []
            for slide_num, slide in enumerate(prs.slides):
                slide_lines = [f"--- Slide {slide_num + 1} ---"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                slide_lines.append(para.text)
                if len(slide_lines) > 1:
                    text_parts.append("\n".join(slide_lines))
            return "\n\n".join(text_parts) if text_parts else None
        except ImportError:
            logger.warning("python-pptx not installed, cannot search PPTX files")
            return None
        except Exception as e:
            logger.warning(f"Failed to extract text from PPTX {file_path}: {e}")
            return None

    return None


def extract_text(file_path: Path, project_dir: Path) -> Path | None:
    """Lazy text extraction with mtime-based caching.

    Checks {project_dir}/.clawmate/cache/text/{hash}_{mtime}.txt.
    If cache exists and mtime matches, returns cache path directly.
    Otherwise extracts text and writes new cache file.

    Args:
        file_path: Absolute path to the source file (PDF/DOCX/etc.)
        project_dir: Absolute path to the project root (contains .clawmate/)

    Returns:
        Path to the cached text file, or None if extraction is not supported.
    """
    try:
        stat = file_path.stat()
    except (OSError, PermissionError):
        return None

    cache_dir = project_dir / ".clawmate" / "cache" / "text"
    phash = _path_hash(file_path)
    cache_key = f"{phash}_{int(stat.st_mtime)}"
    cache_file = cache_dir / f"{cache_key}.txt"

    if cache_file.exists():
        return cache_file

    # Extract text
    text = _extract_by_type(file_path)
    if text is None:
        return None

    cache_dir.mkdir(parents=True, exist_ok=True)
    cache_file.write_text(text, encoding="utf-8", errors="replace")

    # Clean up old caches for the same file (different mtime)
    try:
        for old in cache_dir.glob(f"{phash}_*"):
            if old.name != f"{cache_key}.txt":
                old.unlink()
    except OSError:
        pass

    return cache_file


def _find_projects(root_path: Path, target: Path, max_depth: int,
                   exclude_dir: set, timeout_deadline: float) -> List[Path]:
    """BFS walk from target to discover project directories (those containing .clawmate/).

    Returns a list of project root Paths (immediate children of root_path that have .clawmate/).
    Also returns non-project directories that should get filename-only search.
    """
    import time as _time
    projects = []

    # Check if target itself is a project
    if (target / ".clawmate").is_dir():
        projects.append(target)
        return projects  # target is a project, no need to search subdirectories

    queue: deque = deque([(target, 0)])
    seen_dirs = {target}

    while queue:
        if _time.time() > timeout_deadline:
            break
        dir_path, depth = queue.popleft()
        if depth > max_depth:
            continue

        try:
            entries = sorted(dir_path.iterdir(), key=lambda p: (not p.is_dir(), p.name.lower()))
        except (OSError, PermissionError):
            continue

        for entry in entries:
            if _time.time() > timeout_deadline:
                break
            if not entry.is_dir():
                continue
            if entry.name.startswith(".") and entry.name in exclude_dir:
                continue
            # Skip all hidden directories (not in exclude list)
            if entry.name.startswith("."):
                continue
            if entry.name.lower() in exclude_dir:
                continue

            # Check if this is a project
            if (entry / ".clawmate").is_dir():
                projects.append(entry)
                # Don't recurse into project subdirectories — rg handles that
                continue

            if depth < max_depth and entry not in seen_dirs:
                seen_dirs.add(entry)
                queue.append((entry, depth + 1))

    return projects


def _project_cache_root(root_path: Path, target: Path) -> Path | None:
    """Return the nearest project root whose .clawmate cache owns ``target``."""
    candidate = target
    while True:
        if (candidate / ".clawmate").is_dir():
            return candidate
        if candidate == root_path:
            return None
        candidate = candidate.parent


def search_content(query: str, root_id: str, rel_dir: str = "",
                   context_lines: int = 3, max_depth: int = 8,
                   max_filesize_mb: int = 5, timeout: float = 30.0,
                   exclude_ext: list = None, exclude_dir: list = None) -> Dict:
    """Search file contents using ripgrep within project directories.

    Walk strategy:
    1. BFS from starting directory, discover .clawmate/ markers → project dirs
    2. For each project directory: run rg to search all text files
       (rg handles its own recursive walk internally, respecting .gitignore)
    3. For PDF/Office files in projects: extract text to cache, then rg searches cache
    4. For non-project areas: filename-only matching only

    Args:
        query: Search query string (literal, not regex by default)
        root_id: Root identifier
        rel_dir: Relative directory to start search from
        context_lines: Lines of context around each match
        max_depth: Maximum directory depth for project discovery
        max_filesize_mb: Skip files larger than this
        timeout: Search timeout in seconds
        exclude_ext: File extensions to exclude from content search
        exclude_dir: Directory names to exclude from traversal

    Returns:
        {query, root, dir, total_matches, total_files, searched_files,
         elapsed_ms, results_by_file: [{file, match_count, preview_url, matches, mtime, size}]}
    """
    import time as _time
    import subprocess
    import json as _json
    from urllib.parse import quote

    if not _rg_available():
        raise RuntimeError(
            "ripgrep (rg) is not installed. "
            "Install it with: sudo apt install ripgrep"
        )

    _start = _time.time()
    _deadline = _start + timeout

    root_path, target, safe_dir = safe_path(root_id, rel_dir)
    if not target.exists() or not target.is_dir():
        raise FileNotFoundError("Directory not found")

    exclude_ext_set = set(e.lower().lstrip(".") for e in (exclude_ext or []))
    exclude_dir_set = set(d.lower() for d in (exclude_dir or []))

    # Step 1: Search directly from the selected directory.  A project marker is
    # only needed to locate the shared extraction cache, not to enable rg.
    project_root = _project_cache_root(root_path, target)
    cache_to_original: Dict[str, Path] = {}

    # Step 2: Extract documents from this subtree into its owning project cache.
    _processed_count = 0
    if project_root:
        try:
            for entry in target.rglob("*"):
                if _time.time() > _deadline:
                    break
                if not entry.is_file() or _is_junk_path(entry):
                    continue
                if entry.suffix.lower() in _CONTENT_EXTRACT_EXTS:
                    cache_file = extract_text(entry, project_root)
                    if cache_file:
                        cache_to_original[str(cache_file)] = entry
                    _processed_count += 1
        except (OSError, PermissionError):
            pass

    # Step 3: Search source text plus only this request's extracted cache files.
    search_paths = [str(target), *cache_to_original.keys()]

    rg_bin = _RG_PATH or "rg"
    cmd = [
        rg_bin, "--json", "--no-heading",
        "-i",
        "-C", str(context_lines),
        "--max-depth", str(max_depth),
        "--max-filesize", f"{max_filesize_mb}M",
        "--no-ignore-vcs",  # Don't skip .gitignored files (user wants full search)
    ]

    # Exclude common binary/media directories
    for d in exclude_dir_set:
        cmd.extend(["-g", f"!{d}/**"])

    # Exclude binary/media extensions
    for e in exclude_ext_set:
        cmd.extend(["-g", f"!*.{e}"])

    # Also exclude .clawmate/ itself from search
    cmd.extend(["-g", "!.clawmate/**"])

    # Skip all hidden directories (those starting with '.')
    cmd.extend(["-g", "!.*/**"])

    cmd.append(query)
    cmd.extend(search_paths)

    # Step 4: Run rg
    try:
        result = subprocess.run(
            cmd,
            capture_output=True, text=True,
            timeout=max(5, timeout - (_time.time() - _start)),
        )
    except subprocess.TimeoutExpired:
        result = None
    except FileNotFoundError:
        raise RuntimeError("ripgrep (rg) not found. Install with: sudo apt install ripgrep")

    # Step 5: Parse rg JSON output with context tracking
    matches_by_file: Dict[str, Dict] = {}
    total_matches = 0
    searched_files = set()
    # Per-file context buffer: context lines before a match
    _context_buffer: Dict[str, list] = {}  # file_path -> [(line_number, text)]
    _last_was_match: Dict[str, bool] = {}  # file_path -> True if last entry was a match

    if result and result.stdout:
        for line in result.stdout.strip().split("\n"):
            if not line:
                continue
            try:
                entry = _json.loads(line)
            except _json.JSONDecodeError:
                continue

            etype = entry.get("type")
            data = entry.get("data", {})
            path_data = data.get("path", {})
            file_path = path_data.get("text", "")

            if etype == "begin":
                searched_files.add(file_path)
                _context_buffer[file_path] = []
                _last_was_match[file_path] = False

            elif etype == "context":
                lines_data = data.get("lines", {})
                line_number = data.get("line_number", 0)
                ctx_text = lines_data.get("text", "").rstrip("\n")
                buf = _context_buffer.setdefault(file_path, [])
                buf.append((line_number, ctx_text))
                _last_was_match[file_path] = False

            elif etype == "match":
                lines_data = data.get("lines", {})
                line_number = data.get("line_number", 0)
                line_text = lines_data.get("text", "").rstrip("\n")

                # Convert absolute path to relative
                try:
                    abs_path = Path(file_path)
                    rel_path = str(abs_path.relative_to(root_path))
                except (ValueError, TypeError):
                    rel_path = file_path

                # Get stored context buffer
                buf = _context_buffer.get(file_path, [])
                is_after_context = _last_was_match.get(file_path, False)

                if is_after_context and buf and file_path in matches_by_file:
                    prev_matches = matches_by_file[file_path]["matches"]
                    if prev_matches:
                        prev_matches[-1]["context_after"] = [
                            t for _, t in buf
                        ]

                # Context before for this match (same lines serve as both after
                # for previous match and before for this match)
                context_before = [t for _, t in buf]
                _context_buffer[file_path] = []  # clear buffer for next match

                if file_path not in matches_by_file:
                    matches_by_file[file_path] = {
                        "file": rel_path,
                        "match_count": 0,
                        "matches": [],
                    }
                matches_by_file[file_path]["match_count"] += 1
                matches_by_file[file_path]["matches"].append({
                    "line": line_number,
                    "text": line_text,
                    "context_before": context_before,
                    "context_after": [],
                })
                total_matches += 1
                _last_was_match[file_path] = True

            elif etype == "end":
                # Flush any remaining context as after-context for last match
                buf = _context_buffer.pop(file_path, [])
                if buf and file_path in matches_by_file:
                    prev_matches = matches_by_file[file_path]["matches"]
                    if prev_matches:
                        prev_matches[-1]["context_after"] = [
                            t for _, t in buf
                        ]
                _last_was_match.pop(file_path, None)

    # Step 6: Replace extracted cache file paths with their original files.
    resolved_results = {}
    for abs_path, info in matches_by_file.items():
        if abs_path in cache_to_original:
            original_file = cache_to_original[abs_path]
            rel_path = str(original_file.relative_to(root_path))
            if rel_path in resolved_results:
                resolved_results[rel_path]["match_count"] += info["match_count"]
                resolved_results[rel_path]["matches"].extend(info["matches"])
            else:
                resolved_results[rel_path] = {
                    "file": rel_path,
                    "match_count": info["match_count"],
                    "matches": info["matches"],
                }
        else:
            resolved_results[abs_path] = info

    # Step 7: Build final response
    from service import get_public_base_url as _get_base_url

    results_by_file = []
    # Sort by match count descending
    sorted_results = sorted(
        resolved_results.values(),
        key=lambda r: r["match_count"],
        reverse=True,
    )

    for info in sorted_results:
        rel_path = info["file"]
        preview_url = (
            f"/clawmate/preview.html"
            f"?root={quote(root_id, safe='')}"
            f"&file={quote(rel_path, safe='')}"
        )
        # Gather file mtime & size for display in search results
        file_mtime = 0
        file_size = 0
        try:
            abs_path = root_path / rel_path
            if abs_path.exists():
                st = abs_path.stat()
                file_mtime = int(st.st_mtime)
                file_size = st.st_size
        except (OSError, ValueError):
            pass
        results_by_file.append({
            "file": rel_path,
            "match_count": info["match_count"],
            "preview_url": preview_url,
            "matches": info["matches"],
            "mtime": file_mtime,
            "size": file_size,
        })

    # Step 8: Check extractor availability for notice
    extractors = _check_extractors()
    unavailable = []
    _ext_label = {".pdf": "PDF", ".docx": "DOCX", ".xlsx": "XLSX", ".pptx": "PPTX"}
    for ext, ok in extractors.items():
        if not ok:
            unavailable.append({"ext": ext, "label": _ext_label.get(ext, ext)})

    # Step 9: Build final response
    return {
        "query": query,
        "root": root_id,
        "dir": safe_dir,
        "total_matches": total_matches,
        "total_files": len(results_by_file),
        "searched_files": len(searched_files),
        "elapsed_ms": int((_time.time() - _start) * 1000),
        "results_by_file": results_by_file,
        "projects_found": int(project_root is not None),
        "unavailable_types": unavailable,  # file types skipped due to missing deps
    }
